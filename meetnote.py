"""mp3 -> 전사 -> 위계 요약 -> FigJam 도식(mermaid). CLI 한 방."""
import argparse, base64, json, subprocess, sys, tempfile, zlib
from pathlib import Path

MODEL = "claude-opus-5"

SCHEMA = {
    "type": "object",
    "properties": {
        "title": {"type": "string", "description": "회의 제목"},
        "topics": {
            "type": "array",
            "description": "회의에서 다뤄진 메인 주제들. 실제 논의된 시간 순서대로.",
            "items": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "주제 한 줄 (12자 내외)"},
                    "summary": {"type": "string", "description": "이 주제의 핵심 결론 1~2문장"},
                    "points": {
                        "type": "array",
                        "description": "이 주제에서 나온 개별 의견/결정/질문/할일",
                        "items": {
                            "type": "object",
                            "properties": {
                                "text": {"type": "string", "description": "한 줄 (25자 내외)"},
                                "kind": {"type": "string", "enum": ["opinion", "decision", "question", "action"]},
                                "speaker": {"type": "string", "description": "화자, 모르면 빈 문자열"},
                            },
                            "required": ["text", "kind", "speaker"],
                            "additionalProperties": False,
                        },
                    },
                },
                "required": ["title", "summary", "points"],
                "additionalProperties": False,
            },
        },
        "decisions": {"type": "array", "items": {"type": "string"}, "description": "회의 전체의 확정 결정사항"},
        "action_items": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "task": {"type": "string"},
                    "owner": {"type": "string"},
                    "due": {"type": "string"},
                },
                "required": ["task", "owner", "due"],
                "additionalProperties": False,
            },
        },
    },
    "required": ["title", "topics", "decisions", "action_items"],
    "additionalProperties": False,
}

SYSTEM = """너는 회의록 편집자다. 전사문을 받아 위계와 핵심 중심으로 구조화한다.

원칙:
- 위계: 메인 주제(topics) > 그 안의 개별 발언(points). 잡담/중복/말버릇은 버린다.
- 핵심: 각 주제의 summary는 "무엇이 논의되어 어떻게 결론났는지"만 쓴다. 발언 나열 금지.
- topics는 실제 논의 흐름 순서. 주제가 되돌아왔으면 하나로 합친다.
- points의 kind는 정확히 분류한다: decision(확정), action(누가 할 일), question(미해결), opinion(그 외 주장).
- 전사문에 없는 내용은 절대 만들지 않는다. 화자를 모르면 speaker는 빈 문자열.
- 모든 출력은 한국어."""

KIND_STYLE = {  # mermaid 노드 모양 + 클래스
    "decision": ("{{", "}}", "dec"),
    "action": ("[/", "/]", "act"),
    "question": ("{", "}", "que"),
    "opinion": ("(", ")", "opi"),
}


def shrink(src: Path) -> Path:
    """Whisper 25MB 제한 대비: 16kHz 모노 32kbps로 재인코딩. 2시간 회의도 한 파일에 들어감."""
    # ponytail: 2시간 초과분은 잘림 없이 그냥 실패함. 그때 ffmpeg -f segment로 쪼개고 전사문 이어붙이면 됨.
    out = Path(tempfile.mkdtemp()) / "audio.mp3"
    subprocess.run(
        ["ffmpeg", "-nostdin", "-loglevel", "error", "-y", "-i", str(src),
         "-ac", "1", "-ar", "16000", "-b:a", "32k", str(out)],
        check=True,
    )
    return out


PIN = Path(__file__).parent / "meetnote.app/Contents/MacOS/pin"


def transcribe(src: Path, on_progress=None) -> str:
    """macOS 온디바이스 전사(SpeechAnalyzer). API 키도 비용도 필요 없다.

    앱 번들이 없으면(./pin.sh 를 한 번도 안 돌렸으면) OpenAI Whisper로 넘어간다."""
    if PIN.exists():
        out = src.parent / "transcript.txt"
        p = subprocess.Popen([str(PIN), "--transcribe", str(src), str(out)],
                             stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        for line in p.stdout:
            if line.startswith("P ") and on_progress:
                on_progress(float(line[2:]))
        if p.wait() != 0:
            raise RuntimeError(f"전사 실패: {p.stderr.read().strip()}")
        return out.read_text()

    from openai import OpenAI
    small = shrink(src)
    if small.stat().st_size > 25 * 1024 * 1024:
        sys.exit(f"오디오가 너무 김 ({small.stat().st_size/1e6:.0f}MB > 25MB). 분할 전사가 필요함.")
    with open(small, "rb") as f:
        return OpenAI().audio.transcriptions.create(model="whisper-1", file=f, language="ko").text


def summarize(transcript: str) -> dict:
    import anthropic
    with anthropic.Anthropic().messages.stream(
        model=MODEL,
        max_tokens=32000,
        system=SYSTEM,
        thinking={"type": "adaptive"},
        output_config={"format": {"type": "json_schema", "schema": SCHEMA}},
        messages=[{"role": "user", "content": f"<transcript>\n{transcript}\n</transcript>"}],
    ) as stream:
        msg = stream.get_final_message()
    return json.loads(next(b.text for b in msg.content if b.type == "text"))


def esc(s: str) -> str:
    return s.replace('"', "#quot;").replace("\n", " ").strip()


def to_mermaid(s: dict) -> str:
    """메인 주제는 세로 흐름 화살표, 개별 의견은 주제에서 갈라지는 가지."""
    lines = ["flowchart TD", f'  START(["{esc(s["title"])}"])']
    prev = "START"
    for i, t in enumerate(s["topics"]):
        tid = f"T{i}"
        lines.append(f'  {tid}["<b>{esc(t["title"])}</b><br/>{esc(t["summary"])}"]')
        lines.append(f"  {prev} ==> {tid}")
        for j, p in enumerate(t["points"]):
            o, c, cls = KIND_STYLE[p["kind"]]
            label = esc(p["text"]) + (f' <i>({esc(p["speaker"])})</i>' if p.get("speaker") else "")
            pid = f"{tid}p{j}"
            lines.append(f'  {pid}{o}"{label}"{c}:::{cls}')
            lines.append(f"  {tid} -.-> {pid}")
        prev = tid
    lines += [
        "  classDef dec fill:#d4f4dd,stroke:#2b8a3e",
        "  classDef act fill:#ffe8cc,stroke:#e8590c",
        "  classDef que fill:#ffe3e3,stroke:#c92a2a",
        "  classDef opi fill:#e7f5ff,stroke:#1971c2",
    ]
    return "\n".join(lines)


def to_markdown(s: dict) -> str:
    icon = {"decision": "✅", "action": "🔨", "question": "❓", "opinion": "💬"}
    md = [f"# {s['title']}\n"]
    for i, t in enumerate(s["topics"], 1):
        md.append(f"## {i}. {t['title']}\n\n{t['summary']}\n")
        for p in t["points"]:
            who = f" — {p['speaker']}" if p.get("speaker") else ""
            md.append(f"- {icon[p['kind']]} {p['text']}{who}")
        md.append("")
    if s["decisions"]:
        md.append("## 결정사항\n")
        md += [f"- {d}" for d in s["decisions"]] + [""]
    if s["action_items"]:
        md.append("## 액션 아이템\n\n| 할 일 | 담당 | 기한 |\n|---|---|---|")
        md += [f"| {a['task']} | {a['owner'] or '-'} | {a['due'] or '-'} |" for a in s["action_items"]]
    return "\n".join(md) + "\n"


# ---------- export ----------

LABEL = {"decision": "[결정]", "action": "[할일]", "question": "[질문]", "opinion": "[의견]"}


def mermaid_live_url(mmd: str) -> str:
    """mermaid.live는 pako(zlib) 압축 + base64url을 fragment로 받는다."""
    raw = json.dumps({"code": mmd, "mermaid": {"theme": "default"}}).encode()
    return "https://mermaid.live/edit#pako:" + base64.urlsafe_b64encode(zlib.compress(raw, 9)).decode()


def _docx(s: dict, path: Path):
    from docx import Document
    d = Document()
    d.add_heading(s["title"], 0)
    for i, t in enumerate(s["topics"], 1):
        d.add_heading(f"{i}. {t['title']}", 1)
        d.add_paragraph(t["summary"])
        for pt in t["points"]:
            who = f" — {pt['speaker']}" if pt.get("speaker") else ""
            d.add_paragraph(f"{LABEL[pt['kind']]} {pt['text']}{who}", style="List Bullet")
    if s["decisions"]:
        d.add_heading("결정사항", 1)
        for x in s["decisions"]:
            d.add_paragraph(x, style="List Bullet")
    if s["action_items"]:
        d.add_heading("액션 아이템", 1)
        tbl = d.add_table(rows=1, cols=3)
        tbl.style = "Table Grid"
        for c, h in zip(tbl.rows[0].cells, ("할 일", "담당", "기한")):
            c.text = h
        for a in s["action_items"]:
            r = tbl.add_row().cells
            r[0].text, r[1].text, r[2].text = a["task"], a["owner"] or "-", a["due"] or "-"
    d.save(path)


def _pptx(s: dict, path: Path):
    from pptx import Presentation
    pr = Presentation()
    cover = pr.slides.add_slide(pr.slide_layouts[0])
    cover.shapes.title.text = s["title"]
    cover.placeholders[1].text = " · ".join(t["title"] for t in s["topics"])
    for i, t in enumerate(s["topics"], 1):
        sl = pr.slides.add_slide(pr.slide_layouts[1])
        sl.shapes.title.text = f"{i}. {t['title']}"
        tf = sl.placeholders[1].text_frame
        tf.text = t["summary"]
        for pt in t["points"]:
            par = tf.add_paragraph()
            par.text = f"{LABEL[pt['kind']]} {pt['text']}" + (f" — {pt['speaker']}" if pt.get("speaker") else "")
            par.level = 1
    if s["decisions"] or s["action_items"]:
        sl = pr.slides.add_slide(pr.slide_layouts[1])
        sl.shapes.title.text = "결정사항 · 액션 아이템"
        tf = sl.placeholders[1].text_frame
        tf.text = "결정사항"
        for x in s["decisions"]:
            par = tf.add_paragraph(); par.text = x; par.level = 1
        par = tf.add_paragraph(); par.text = "액션 아이템"
        for a in s["action_items"]:
            par = tf.add_paragraph()
            par.text = f"{a['task']} — {a['owner'] or '-'} ({a['due'] or '-'})"
            par.level = 1
    pr.save(path)


def export(s: dict, out: Path, target: str) -> str:
    """target별 산출물을 만들고, 데스크탑에서 열 대상(파일 경로 또는 URL)을 돌려준다."""
    out.mkdir(parents=True, exist_ok=True)
    (out / "summary.json").write_text(json.dumps(s, ensure_ascii=False, indent=2))
    (out / "summary.md").write_text(to_markdown(s))
    if target == "figjam":
        mmd = to_mermaid(s)
        (out / "diagram.mmd").write_text(mmd)
        # ponytail: FigJam 보드 직접 생성은 Figma 플러그인이나 MCP가 필요하다.
        # 지금은 mermaid.live를 열어 SVG/PNG로 받아 FigJam에 붙이는 경로.
        return mermaid_live_url(mmd)
    if target == "word":
        _docx(s, out / "summary.docx")
        return str(out / "summary.docx")
    if target == "ppt":
        _pptx(s, out / "summary.pptx")
        return str(out / "summary.pptx")
    raise ValueError(f"모르는 target: {target}")


def main():
    ap = argparse.ArgumentParser(description="회의 mp3 -> 위계 요약 + FigJam 도식")
    ap.add_argument("audio", type=Path, help="mp3/m4a/wav 등 ffmpeg가 읽는 오디오")
    ap.add_argument("-o", "--out", type=Path, default=Path("out"))
    ap.add_argument("--transcript", type=Path, help="전사 건너뛰고 이 텍스트 파일 사용")
    ap.add_argument("-t", "--target", choices=["figjam", "word", "ppt"], help="export 대상")
    ap.add_argument("--open", action="store_true", help="결과물을 바로 열기")
    a = ap.parse_args()

    a.out.mkdir(parents=True, exist_ok=True)
    if a.transcript:
        text = a.transcript.read_text()
    else:
        print("전사 중...", file=sys.stderr)
        text = transcribe(a.audio)
        (a.out / "transcript.txt").write_text(text)

    print("요약 중...", file=sys.stderr)
    s = summarize(text)
    dest = export(s, a.out, a.target or "figjam")
    print(f"완료 -> {a.out}/  열 대상: {dest}")
    if a.open:
        subprocess.run(["open", dest])


if __name__ == "__main__":
    main()
